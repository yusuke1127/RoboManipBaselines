import argparse
import math
import os
import re
import subprocess
import sys
import tempfile
from abc import ABC, abstractmethod
from datetime import datetime
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from ruamel.yaml import YAML
from tqdm import tqdm

import cv2
import imageio
import matplotlib.pyplot as plt
import numpy as np

BAR_CHART_TOTAL_WIDTH = 0.8
BAR_CHART_PLOT_CAPSIZE = 5

GREEN = (0, 128, 0)
RED = (0, 0, 255)
WHITE = (255, 255, 255)
BLACK = (0, 0, 0)
GRAY = (128, 128, 128)

CAPTURE_TRIM_COLOR_THRESHOLD = 128


def parse_args():
    parser = argparse.ArgumentParser()

    parser.add_argument(
        "--max_video_width",
        "-w",
        type=int,
        default=1920,
        help="maximum width to which the video will be scaled down if it is too large",
    )
    parser.add_argument("--max_agenda_lines", type=int, default=13)
    parser.add_argument("--common_title_font_size_pt", type=float, default=24)
    parser.add_argument("--description_font_size_pt", type=float, default=18)
    parser.add_argument(
        "--plot_outline_position_and_size_in_inches",
        nargs=4,
        default=[1.57, 2.52, None, 4.92],
        help=(
            "specify the position and size of the plot summary of summaries "
            "in inches: "
            "left, top, width, height"
        ),
    )
    parser.add_argument(
        "--plot_summary_position_and_size_in_inches",
        nargs=4,
        default=[0.25, 3.6, 4.4, None],
        help=(
            "specify the position and size of the plot image in inches: "
            "left, top, width, height"
        ),
    )
    parser.add_argument(
        "--captured_start_end_img_pos_and_size_in_inches",
        nargs=4,
        default=[4.34, 4.25, 2.55, None],
        help=(
            "specify the position and size of the captured start/end images "
            "in inches: "
            "left, top, width, height"
        ),
    )
    parser.add_argument(
        "--video_position_and_size_in_inches",
        nargs=4,
        default=[0.15, 1.25, 9.75, 5.75],
        help=(
            "specify the position and size of video in inches: "
            "left, top, width, height"
        ),
    )
    parser.add_argument(
        "--video_caption_position_and_size_in_inches",
        nargs=4,
        default=[4, 7, 2, 2],
        help=(
            "specify the position and size of video caption in inches: "
            "left, top, width, height"
        ),
    )
    parser.add_argument("--video_column_num", default=3)
    parser.add_argument(
        "--jump_position_and_size_in_inches",
        nargs=4,
        default=[0.18, 0, 2, 0.33],
        help=(
            "specify the position and size of jump in inches: "
            "left, top, width, height"
        ),
    )
    parser.add_argument("--jump_font_size_pt", type=float, default=14)
    parser.add_argument(
        "--pagenum_position_and_size_in_inches",
        nargs=4,
        default=[9, 0, 0.9, 0.3],
        help=(
            "specify the position and size of page number in inches: "
            "left, top, width, height"
        ),
    )
    parser.add_argument("--pagenum_font_size_pt", type=float, default=12)
    parser.add_argument(
        "--xlabel_rotation", type=float, dest="xlabel_rotation", default=10
    )

    curr_dir = os.path.dirname(os.path.abspath(__file__))
    parser.add_argument(
        "-i",
        "--in_yaml_filename",
        default=os.path.join(curr_dir, "..", "report", "sample.yaml"),
    )
    parser.add_argument(
        "-o",
        "--out_pptx_filename",
        default=os.path.join(curr_dir, ".", "output.pptx"),
    )

    parsed_args = parser.parse_args()
    tqdm.write(f"{parsed_args=}")
    return parsed_args


class PosSize:
    def __init__(self, left, top, width, height):
        self.left = left
        self.top = top
        self.width = width
        self.height = height

    def as_tuple(self):
        return (self.left, self.top, self.width, self.height)


class AbstractPlotSlideinfo(ABC):
    """
    abstract class for plotting information about successes onto slides with
    title and content layout
    """

    def __init__(self, yaml_idx):
        self.yaml_idx = yaml_idx

        self.title_str = ""
        self.successes = {}
        self.xlabels = []
        self.policy_names = []

    def get_video_root(self, media_dir, xlabel, policy_name):
        return os.path.join(
            media_dir,
            os.path.splitext(self.yaml_idx)[0],
            f"{xlabel}_{policy_name}",
        )

    def lookup(self, c_dict):
        return c_dict[self.__class__.__name__]

    @abstractmethod
    def generate_description(self):
        raise NotImplementedError()


class OutlineSlideinfo(AbstractPlotSlideinfo):
    """summary of summaries slide information"""

    def __init__(self, yaml_idx):
        super().__init__(yaml_idx)
        self.description = ""

    def generate_description(self):
        return self.description


class SummarySlideinfo(AbstractPlotSlideinfo):
    def __init__(self, yaml_idx):
        super().__init__(yaml_idx)
        self.headline = ""  # <Env name>_<Val name>_<date>
        self.text_content_dict = {}  # robot, env, task, validation

    def generate_description(self):
        """description: robot, env, task, validation"""
        desc_texts = []
        robot_in_env = " in ".join(
            [
                s
                for s in [
                    self.text_content_dict.get("robot", ""),
                    self.text_content_dict.get("env", ""),
                ]
                if s
            ]
        )
        if robot_in_env:
            desc_texts.append(f"環境：{robot_in_env}")
        if "task" in self.text_content_dict:
            desc_texts.append(f"タスク：{self.text_content_dict['task']}")
        if "validation" in self.text_content_dict:
            desc_texts.append("検証内容：" + self.text_content_dict["validation"])
        return "\n".join(desc_texts)


def check_value(condition_val, optional_msg=None):
    if condition_val:
        return True

    msg = "\tWarn: Unexpected value detected."
    if optional_msg:
        msg += f"\n\t\t{optional_msg}"
    tqdm.write(msg, file=sys.stderr)

    return False


def adjust_size(
    curre_width, curre_height, valid_left, valid_top, valid_width, valid_height
):
    curre_aspect_ratio = curre_width / curre_height

    valid_center_x = valid_left + valid_width / 2
    valid_center_y = valid_top + valid_height / 2

    adjed_width = valid_width
    adjed_height = valid_width / curre_aspect_ratio

    if adjed_height > valid_height:
        adjed_height = valid_height
        adjed_width = adjed_height * curre_aspect_ratio

    if adjed_width > valid_width:
        adjed_width = valid_width
        adjed_height = adjed_width / curre_aspect_ratio

    adjed_left = valid_center_x - adjed_width / 2
    adjed_top = valid_center_y - adjed_height / 2

    return adjed_left, adjed_top, adjed_width, adjed_height


def read_start_end_frame(video_path):
    cap = cv2.VideoCapture(video_path)
    assert check_value(cap.isOpened()), f"Could not open video. [{video_path=}]"
    ret, start_frame = cap.read()
    assert check_value(ret), "Could not read frame."
    cap.set(  # go to the end frame
        cv2.CAP_PROP_POS_FRAMES, cap.get(cv2.CAP_PROP_FRAME_COUNT) - 1
    )
    ret, end_frame = cap.read()
    assert check_value(ret), "Could not read frame."
    cap.release()

    return start_frame, end_frame


class PresentationHandler:
    def __init__(
        self,
        in_yaml_filename,
        out_pptx_filename,
        max_video_width,
        max_agenda_lines,
        common_title_font_size_pt,
        description_font_size_pt,
        plot_outline_pos_size,
        plot_summary_pos_size,
        captu_img_pos_size,
        video_pos_size,
        vcaption_pos_size,
        video_column_num,
        jump_pos_size,
        jump_font_size_pt,
        pagenum_pos_size,
        pagenum_font_size_pt,
        xlabel_rotation,
    ):
        self.in_yaml_filename = in_yaml_filename
        self.out_pptx_filename = out_pptx_filename
        self.media_dir = os.path.dirname(in_yaml_filename)
        self.max_video_width = max_video_width
        self.max_agenda_lines = max_agenda_lines
        self.common_title_font_size_pt = common_title_font_size_pt
        self.description_font_size_pt = description_font_size_pt
        self.plot_pos_sizes = {
            OutlineSlideinfo.__name__: plot_outline_pos_size,
            SummarySlideinfo.__name__: plot_summary_pos_size,
        }
        self.captu_img_pos_size = captu_img_pos_size
        self.video_pos_size = video_pos_size
        self.vcaption_pos_size = vcaption_pos_size
        self.video_column_num = video_column_num
        self.jump_pos_size = jump_pos_size
        self.jump_font_size_pt = jump_font_size_pt
        self.pagenum_pos_size = pagenum_pos_size
        self.pagenum_font_size_pt = pagenum_font_size_pt
        self.xlabel_rotation = xlabel_rotation

        self.presentation = Presentation()
        self.slideinfo_list = []

    def save_plotted_bar(self, xlabels, policy_names, successes, out_fig_filename):
        w_bar = BAR_CHART_TOTAL_WIDTH / len(policy_names)
        x_pos = np.arange(len(xlabels))

        plt.ylabel("Accuracy rate")
        plt.xticks(
            x_pos + w_bar * (len(policy_names) - 1) / 2,
            xlabels if len(xlabels) >= 2 else [""] * len(xlabels),
            rotation=self.xlabel_rotation,
            ha="right",
        )
        plt.ylim([0.0, 1.0])

        for i, policy_name in enumerate(policy_names):
            plt.bar(
                x_pos + i * w_bar,
                height=[np.mean(successes[xlabel][policy_name]) for xlabel in xlabels],
                width=w_bar,
                yerr=[
                    np.std(np.mean(successes[xlabel][policy_name], axis=1))
                    for xlabel in xlabels
                ],
                label=policy_name,
                capsize=BAR_CHART_PLOT_CAPSIZE,
            )
        plt.legend()

        plt.savefig(out_fig_filename)
        plt.clf()
        plt.close()

    def view_result(self, rdata):
        assert isinstance(rdata, dict), f"{type(rdata)=}"
        xlabels = []
        policy_names = []
        successes = {}
        for xlabel, policies in rdata.items():
            assert isinstance(xlabel, str), f"{type(xlabel)=}"
            assert isinstance(policies, dict), f"{type(policies)=}"
            if xlabel not in xlabels:
                xlabels.append(xlabel)
            if xlabel not in successes:
                successes[xlabel] = {}
            for policy_name, policy_val in policies.items():
                assert isinstance(policy_name, str), f"{type(policy_name)=}"
                assert isinstance(policy_val, list), f"{type(policy_val)=}"
                if policy_name not in policy_names:
                    policy_names.append(policy_name)
                if policy_name not in successes[xlabel]:
                    successes[xlabel][policy_name] = []
                for results in policy_val:
                    assert isinstance(results, dict), f"{type(results)=}"
                    for result_name, result_val in results.items():
                        assert isinstance(result_name, str), f"{type(result_name)=}"
                        if result_name == "success":
                            assert isinstance(result_val, list), f"{type(result_val)=}"
                            successes[xlabel][policy_name].append(result_val)
                        elif result_name == "video_url":
                            # internal only information
                            assert isinstance(result_val, str), f"{type(result_val)=}"
                        else:
                            raise AssertionError(f"{result_name=}")
        return xlabels, policy_names, successes

    def layout_of(self, layout_name):
        for layout in self.presentation.slide_layouts:
            if layout.name == layout_name:
                return layout

        layout_names = [s.name for s in self.presentation.slide_layouts]
        raise AssertionError(f"{(layout_name, layout_names)}=")

    def select_video_tile(self, slideinfo):
        xlabel = sorted(slideinfo.xlabels)[-1]
        best_accuracy = sorted(
            [
                (np.mean(succ), poli)
                for poli, succ in slideinfo.successes[xlabel].items()
            ]
        )[-1]
        best_pol = best_accuracy[-1]
        video_root = slideinfo.get_video_root(self.media_dir, xlabel, best_pol)
        suc = slideinfo.successes[xlabel][best_pol]
        best_suc_ind = np.argmax(suc)
        ind_row = best_suc_ind // self.video_column_num
        ind_col = best_suc_ind % self.video_column_num
        video_row_num = len(suc[0]) // self.video_column_num

        return video_root, ind_row, ind_col, video_row_num

    def add_captu_imgs_to_slide(self, slideinfo, summary_id_num):
        """add captured start end images to slide"""
        slide = self.presentation.slides[summary_id_num]

        # captured start end images
        video_root, ind_row, ind_col, video_row_num = self.select_video_tile(slideinfo)
        if not video_root:
            return
        start_frame, end_frame = read_start_end_frame(
            os.path.join(self.media_dir, video_root + ".wmv")
        )
        vtile_h = start_frame.shape[0] // video_row_num
        vtile_w = start_frame.shape[1] // self.video_column_num
        ilef, itop, iwid, ihei = self.captu_img_pos_size.as_tuple()
        for n, frame in enumerate([start_frame, end_frame]):
            tile_img = frame[
                vtile_h * ind_row : vtile_h * (ind_row + 1),
                vtile_w * ind_col : vtile_w * (ind_col + 1),
            ]
            for trim_color in [GREEN, RED, WHITE]:
                dist_mask = np.where(
                    np.mean(np.abs(tile_img - trim_color), axis=2)
                    > CAPTURE_TRIM_COLOR_THRESHOLD
                )
                inds = [func(dist_mask[ax]) for ax in [0, 1] for func in [min, max]]
                tile_img = tile_img[inds[0] : inds[1], inds[2] : inds[3]]
            _, captu_img_path = tempfile.mkstemp(suffix=".png")
            cv2.imwrite(captu_img_path, tile_img)
            with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp_file:
                Image.fromarray(imageio.get_reader(captu_img_path).get_data(0)).save(
                    tmp_file.name
                )
                slide.shapes.add_picture(
                    captu_img_path, ilef + (iwid * 1.2 * n), itop, iwid, ihei
                )

        # arrow
        iwid = iwid or slide.shapes[-1].width
        ihei = ihei or slide.shapes[-1].height
        right_arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            ilef + iwid * 1.04,
            itop + ihei / 2.2,
            iwid / 8.5,
            ihei / 5.5,
        )
        right_arrow.fill.background()  # transparent
        right_arrow.line.color.rgb = RGBColor(*BLACK)
        right_arrow.line.width = Pt(2)
        right_arrow.shadow.inherit = False

    def convert_mp4_to_wmv(self, slideinfo):
        ext2path_list = []
        for xlabel in slideinfo.xlabels:
            for policy_name in slideinfo.policy_names:
                ext2path = {
                    ext: slideinfo.get_video_root(self.media_dir, xlabel, policy_name)
                    + f".{ext}"
                    for ext in ["mp4", "wmv"]
                }
                if os.path.exists(ext2path["wmv"]):
                    continue
                assert os.path.exists(ext2path["mp4"]), f"{ext2path['mp4']=}"
                ext2path_list.append(ext2path)

        if not ext2path_list:
            return

        pbar = tqdm(ext2path_list, desc=self.convert_mp4_to_wmv.__name__, leave=False)
        for ext2path in pbar:
            wmv_dir = os.path.dirname(ext2path["wmv"])
            if not os.path.exists(wmv_dir):
                os.makedirs(wmv_dir)
            command = [
                "ffmpeg",
                "-i",
                ext2path["mp4"],
                "-vf",
                "scale="
                + str(self.max_video_width)
                + ":-2:"
                + "force_original_aspect_ratio="
                + "decrease",
                ext2path["wmv"],
                "-loglevel",
                "warning",
            ]
            pbar.leave = True
            tqdm.write(f"{command=}")
            subprocess.run(command, shell=False, check=True)

    def add_agenda_content(self, new_title_str):
        for id_num_slide in range(1, self.num_agenda_slides + 1):
            p_holder = self.presentation.slides[id_num_slide].placeholders[1]

            n_lines = p_holder.text.count("\n") + 1
            if n_lines >= self.max_agenda_lines:
                continue  # next agenda slide

            cat_str = f"{p_holder.text}\n{new_title_str}"
            p_holder.text = cat_str.lstrip("\n")
            return

        raise AssertionError()

    def add_video_to_slide(
        self,
        video_count,
        video_root,
        xlabel,
        policy_name,
        len_xlabels,
        len_policy_names,
        title_str,
    ):
        # title
        video_path = os.path.join(self.media_dir, video_root) + ".wmv"
        new_video_slide = self.presentation.slides.add_slide(
            self.layout_of("Title Only")
        )
        new_video_slide.shapes.title.text = (
            title_str
            + "（"
            + str(video_count)
            + "/"
            + str(len_xlabels * len_policy_names)
            + "）"
        )

        # textbox
        paragraphs = new_video_slide.shapes.add_textbox(
            *self.vcaption_pos_size.as_tuple()
        ).text_frame.paragraphs
        paragraphs[0].add_run().text = (
            ": ".join([xlabel, policy_name]) if len_xlabels >= 2 else policy_name
        )
        for paragraph in paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER

        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp_file:
            # save temporary image file
            Image.fromarray(imageio.get_reader(video_path).get_data(0)).save(
                tmp_file.name
            )

            # video capture
            capture = cv2.VideoCapture(video_path)
            cwid = int(capture.get(cv2.CAP_PROP_FRAME_WIDTH))
            chei = int(capture.get(cv2.CAP_PROP_FRAME_HEIGHT))
            capture.release()

            # add_movie
            movie = new_video_slide.shapes.add_movie(
                video_path,
                *adjust_size(cwid, chei, *self.video_pos_size.as_tuple()),
                poster_frame_image=tmp_file.name,
                mime_type="video/mp4",
            )
            parent_elem = movie._element.getparent().getparent()
            for elem in parent_elem.getnext().getnext().iter():
                if "cond" in elem.tag:
                    elem.set("delay", "0")

    def add_videos_to_slide(self, slideinfo):
        video_id_num_list = []
        video_count = 0
        for xlabel in tqdm(
            slideinfo.xlabels,
            leave=len(slideinfo.xlabels) > 5,
            desc=self.add_video_to_slide.__name__,
        ):
            for policy_name in slideinfo.policy_names:
                video_root = slideinfo.get_video_root(
                    self.media_dir, xlabel, policy_name
                )
                video_count += 1
                self.add_video_to_slide(
                    video_count,
                    video_root,
                    xlabel,
                    policy_name,
                    len(slideinfo.xlabels),
                    len(slideinfo.policy_names),
                    slideinfo.title_str,
                )
                video_id_num_list.append(len(self.presentation.slides) - 1)
        return video_id_num_list

    def add_jump_to_slide(self, summary_id_num, video_id_num_list):
        if not video_id_num_list:
            return

        def set_font(font):
            font.size = Pt(self.jump_font_size_pt)
            font.color.rgb = RGBColor(*GRAY)

        # to video slide
        shape = self.presentation.slides[summary_id_num].shapes.add_textbox(
            *self.jump_pos_size.as_tuple()
        )
        run = shape.text_frame.paragraphs[0].add_run()
        set_font(run.font)
        run.text = "▽動画スライドにジャンプ"
        shape.click_action.target_slide = self.presentation.slides[
            min(video_id_num_list)
        ]

        # to summary slide
        for video_id_num in video_id_num_list:
            shape = self.presentation.slides[video_id_num].shapes.add_textbox(
                *self.jump_pos_size.as_tuple()
            )
            run = shape.text_frame.paragraphs[0].add_run()
            set_font(run.font)
            run.text = "△サマリースライドに戻る"
            shape.click_action.target_slide = self.presentation.slides[summary_id_num]

    def follow_index_for_append(self, index, new_slideinfo):
        assert isinstance(index, dict), f"{type(index)=}"
        assert all((k in ["label", "path"]) for k in index.keys()), f"{index.keys()=}"
        label = ""
        if "label" in index:
            label = index["label"]
        if not label:
            task_match = re.search(
                r"/([^/]*?)_[0-9]+", index["path"], flags=re.IGNORECASE
            )
            assert task_match, f"{index=}"
            label = task_match.group(1)[:16]
        with open(
            os.path.join(os.path.dirname(self.in_yaml_filename), index["path"]),
            "r",
            encoding="utf-8",
        ) as rfile:
            loaded_yaml = YAML().load(rfile)
            result_list = [
                return_value
                for yaml_dict in loaded_yaml.values()
                for key, return_value in yaml_dict.items()
                if key == "results"
            ]
            assert len(result_list) == 1, f"{len(result_list)=}"
            _, new_pol_names, new_successes = self.view_result(result_list[0])
            if label not in new_slideinfo.xlabels:
                new_slideinfo.xlabels.append(label)
            if label not in new_slideinfo.successes:
                new_slideinfo.successes[label] = {}
            for new_pol_name in new_pol_names:
                assert isinstance(new_pol_name, str), f"{type(new_pol_name)=}"
                if new_pol_name not in new_slideinfo.policy_names:
                    new_slideinfo.policy_names.append(new_pol_name)
                if new_pol_name not in new_slideinfo.successes[label]:
                    new_slideinfo.successes[label][new_pol_name] = []
                for new_successes_value in new_successes.values():
                    assert isinstance(
                        new_successes_value, dict
                    ), f"{type(new_successes_value)=}"
                    result_vals = new_successes_value[new_pol_name]
                    assert isinstance(result_vals, list), f"{type(result_vals)=}"
                    for result_val in result_vals:
                        assert isinstance(result_val, list), f"{type(result_val)=}"
                        new_slideinfo.successes[label][new_pol_name].append(result_val)

    def parse_yaml_recursively(self, rdata, yaml_idx):
        if list(rdata.keys())[0] == "Index":
            for yaml_idxs in rdata.values():
                assert isinstance(yaml_idxs, list), f"{type(yaml_idxs)=}"
                for ym_idx in yaml_idxs:
                    assert isinstance(ym_idx, str), f"{type(ym_idx)=}"
                    with open(
                        os.path.join(os.path.dirname(self.in_yaml_filename), ym_idx),
                        "r",
                        encoding="utf-8",
                    ) as rfile:
                        self.parse_yaml_recursively(YAML().load(rfile), ym_idx)
            return
        if list(rdata.keys())[0] == "Outline":
            assert len(rdata.keys()) == 1, f"{rdata.keys()=}"
            assert isinstance(rdata, dict), f"{type(rdata)=}"
            tqdm.write("- Outline")
            new_outline_slideinfo = OutlineSlideinfo(yaml_idx)
            for item_type, item_val in rdata["Outline"].items():
                if item_type == "title":
                    assert isinstance(item_val, str), f"{type(item_val)=}"
                    assert (
                        not new_outline_slideinfo.title_str
                    ), f"{new_outline_slideinfo.title_str=}"
                    new_outline_slideinfo.title_str = item_val
                    continue
                if item_type == "description":
                    assert isinstance(item_val, str), f"{type(item_val)=}"
                    assert (
                        not new_outline_slideinfo.description
                    ), f"{new_outline_slideinfo.description=}"
                    new_outline_slideinfo.description = item_val
                    continue
                if item_type == "index":
                    assert isinstance(item_val, list), f"{type(item_val)=}"
                    for index in item_val:
                        self.follow_index_for_append(index, new_outline_slideinfo)
                    continue
                raise AssertionError(f"{item_type=}")
            self.slideinfo_list.append(new_outline_slideinfo)
            return
        new_summary_slideinfo = SummarySlideinfo(yaml_idx)
        if not new_summary_slideinfo.headline:
            assert isinstance(rdata, dict), f"{type(rdata)=}"
            for headline, rdata_val in rdata.items():
                tqdm.write(f"- {headline}")
                new_summary_slideinfo.headline = headline
                for key, val in rdata_val.items():
                    if key == "title":
                        assert isinstance(val, str), f"{type(val)=}"
                        assert not new_summary_slideinfo.title_str
                        new_summary_slideinfo.title_str = val
                        continue
                    if key in ["robot", "env", "task", "validation"]:
                        assert isinstance(val, str), f"{type(val)=}"
                        assert (
                            key not in new_summary_slideinfo.text_content_dict
                        ), f"{key=}"
                        new_summary_slideinfo.text_content_dict[key] = val
                        continue
                    if key == "results":
                        xlabels, pol_names, successes = self.view_result(val)
                        assert not new_summary_slideinfo.xlabels
                        assert not new_summary_slideinfo.policy_names
                        assert not new_summary_slideinfo.successes
                        new_summary_slideinfo.xlabels = xlabels
                        new_summary_slideinfo.policy_names = pol_names
                        new_summary_slideinfo.successes = successes
                        continue
                    raise AssertionError(f"{key=}")
            self.slideinfo_list.append(new_summary_slideinfo)
            return
        raise AssertionError(f"{(rdata, type(rdata))=}")

    def create_slide(self):
        def create_opening_title_slide():
            for _ in tqdm([1], desc=f"{create_opening_title_slide.__name__}"):
                # slide
                opening_slide = self.presentation.slides.add_slide(
                    self.layout_of("Title Slide")
                )
                # title
                opening_slide.shapes.title.text = (
                    "Benchmark results of RoboManipBaselines"
                )
                # subtitle
                opening_slide.placeholders[1].text = "\n".join(
                    [
                        "https://github.com/isri-aist/RoboManipBaselines",
                        datetime.now().strftime("%Y/%m/%d"),
                    ]
                )
                placeholders = opening_slide.placeholders
                for paragraph in placeholders[1].text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(self.common_title_font_size_pt)
                        run.font.color.rgb = RGBColor(*BLACK)

        def create_agenda_slide():
            self.num_agenda_slides = math.ceil(
                len(self.slideinfo_list) / self.max_agenda_lines
            )
            for i in tqdm(
                range(1, self.num_agenda_slides + 1),
                desc=f"{create_agenda_slide.__name__}",
            ):
                # slide
                agenda_slide = self.presentation.slides.add_slide(
                    self.layout_of("Title and Content")
                )
                # title
                agenda_slide.shapes.title.text = "目次"
                if self.num_agenda_slides >= 2:
                    agenda_slide.shapes.title.text += (
                        "（" + f"{i}/{self.num_agenda_slides}" + "）"
                    )
                # content
                agenda_slide.placeholders[1].text = ""

        def create_plot_slide():
            """outline or summary"""
            for new_slide_info in tqdm(
                self.slideinfo_list, desc=f"{create_plot_slide.__name__}"
            ):
                # slide
                new_slide_entity = self.presentation.slides.add_slide(
                    self.layout_of("Title and Content")
                )
                # title
                new_slide_entity.shapes.title.text = new_slide_info.title_str
                # agenda content
                self.add_agenda_content(new_slide_info.title_str)
                # plot
                if check_value(
                    new_slide_info.policy_names,
                    f"{new_slide_info.policy_names=}",
                ):
                    with tempfile.NamedTemporaryFile(
                        suffix=".png", delete=False
                    ) as tfile:
                        self.save_plotted_bar(
                            new_slide_info.xlabels,
                            new_slide_info.policy_names,
                            new_slide_info.successes,
                            tfile.name,
                        )
                        new_slide_entity.shapes.add_picture(
                            tfile.name,
                            *new_slide_info.lookup(self.plot_pos_sizes).as_tuple(),
                        )

                # description
                placeholder = new_slide_entity.placeholders[1]
                placeholder.text = new_slide_info.generate_description()
                # font size
                for paragraph in placeholder.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(self.description_font_size_pt)

        def create_video_slide():
            for summ_id, slideinfo in enumerate(
                tqdm(self.slideinfo_list, desc=f"{create_video_slide.__name__}"),
                start=len(self.presentation.slides) - len(self.slideinfo_list),
            ):
                if isinstance(slideinfo, OutlineSlideinfo):
                    continue  # video is not needed
                assert isinstance(slideinfo, SummarySlideinfo)
                tqdm.write(f"- {slideinfo.headline}")
                if (not slideinfo.xlabels) or (not slideinfo.policy_names):
                    continue  # cannot find video

                # find videos that match xlabel and policy_name
                self.convert_mp4_to_wmv(slideinfo)

                # capture start end images
                self.add_captu_imgs_to_slide(slideinfo, summ_id)

                # add videos
                jump_vid_id_list = self.add_videos_to_slide(slideinfo)

                # add jump
                self.add_jump_to_slide(summ_id, jump_vid_id_list)

        create_opening_title_slide()
        create_agenda_slide()
        create_plot_slide()  # outlines or summarise
        create_video_slide()

        # common: page number
        for i, slide in enumerate(self.presentation.slides, start=1):
            shape = slide.shapes.add_textbox(*self.pagenum_pos_size.as_tuple())
            shape.text_frame.text = f"{i}/{len(self.presentation.slides)}"
            for paragraph in shape.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.RIGHT
                for run in paragraph.runs:
                    run.font.size = Pt(self.pagenum_font_size_pt)

        # common: font size
        for slide in self.presentation.slides:
            # title
            for paragraph in slide.shapes.title.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(self.common_title_font_size_pt)
            # agenda content
            for i_slide in range(1, self.num_agenda_slides + 1):
                for paragraph in (
                    self.presentation.slides[i_slide]
                    .placeholders[1]
                    .text_frame.paragraphs
                ):
                    for run in paragraph.runs:
                        run.font.size = Pt(self.description_font_size_pt)


if __name__ == "__main__":
    try:
        with tqdm(total=100, desc=__name__.strip("_")) as prog_bar:
            args = parse_args()

            tqdm.write(f"# {__name__.strip('_')}")

            tqdm.write(f"## {Presentation.__init__.__name__.strip('_')}")
            p = PresentationHandler(
                args.in_yaml_filename,
                args.out_pptx_filename,
                args.max_video_width,
                args.max_agenda_lines,
                args.common_title_font_size_pt,
                args.description_font_size_pt,
                PosSize(
                    *[
                        Inches(a) if a else a
                        for a in args.plot_outline_position_and_size_in_inches
                    ]
                ),
                PosSize(
                    *[
                        Inches(a) if a else a
                        for a in args.plot_summary_position_and_size_in_inches
                    ]
                ),
                PosSize(
                    *[
                        Inches(a) if a else a
                        for a in args.captured_start_end_img_pos_and_size_in_inches
                    ]
                ),
                PosSize(
                    *[
                        Inches(a) if a else a
                        for a in args.video_position_and_size_in_inches
                    ]
                ),
                PosSize(
                    *[
                        Inches(a) if a else a
                        for a in args.video_caption_position_and_size_in_inches
                    ]
                ),
                args.video_column_num,
                PosSize(
                    *[
                        Inches(a) if a else a
                        for a in args.jump_position_and_size_in_inches
                    ]
                ),
                args.jump_font_size_pt,
                PosSize(
                    *[
                        Inches(a) if a else a
                        for a in args.pagenum_position_and_size_in_inches
                    ]
                ),
                args.pagenum_font_size_pt,
                args.xlabel_rotation,
            )
            prog_bar.update(0.01)

            tqdm.write(f"## {YAML.load.__name__}")
            with open(p.in_yaml_filename, "r", encoding="utf-8") as rtextio:
                p.parse_yaml_recursively(YAML().load(rtextio), p.in_yaml_filename)
            prog_bar.update(0.01)

            tqdm.write(f"## {Presentation.__module__}")
            p.create_slide()
            prog_bar.update(49.98)

            tqdm.write(f"## {p.presentation.save.__name__}")
            p.presentation.save(args.out_pptx_filename)
            prog_bar.update(50.00)

            tqdm.write(f"{args.out_pptx_filename=}")
            tqdm.write("Done.")
    except Exception:
        prog_bar.leave = False
        import traceback

        traceback.print_exc()
